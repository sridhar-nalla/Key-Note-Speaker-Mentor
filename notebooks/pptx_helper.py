# tools/pptx_helper.py
from pptx import Presentation
from pptx.util import Inches, Pt


def create_basic_presentation(slides: list, out_path: str = 'output.pptx'):
    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        title.text = s.get('title','')
        tf = body.text_frame
        for bullet in s.get('bullets', []):
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        # speaker notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = s.get('notes','')
    prs.save(out_path)
    return out_path


def add_or_update_speaker_notes(pptx_path: str, notes_per_slide: dict, out_path: str = None):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        text = notes_per_slide.get(i, '')
        slide.notes_slide.notes_text_frame.text = text
    out = out_path or pptx_path
    prs.save(out)
    return out