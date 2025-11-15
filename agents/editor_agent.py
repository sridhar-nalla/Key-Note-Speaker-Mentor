# agents/editor_agent.py
from tools.pptx_helper import add_or_update_speaker_notes

class EditorAgent:
    def __init__(self, llm=None):
        self.llm = llm

    def generate_notes_for_ppt(self, pptx_path: str):
        # dummy: create simple notes based on slide titles
        from pptx import Presentation
        prs = Presentation(pptx_path)
        notes = {}
        for i, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else f'Slide {i+1}'
            notes[i] = f'Notes for {title} — expand on key points.'
        return notes

    def add_speaker_notes(self, pptx_path: str, out_path: str = None):
        notes = self.generate_notes_for_ppt(pptx_path)
        return add_or_update_speaker_notes(pptx_path, notes, out_path=out_path)